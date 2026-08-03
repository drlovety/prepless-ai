"""
builders.py — PPTX + DOCX builders using python-pptx and python-docx
Maps the website's JSON schema to proper file output.
"""
import json
import os
import re
import tempfile
from pathlib import Path
from typing import Dict, List, Optional

import requests
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches
from pptx.dml.color import RGBColor as PptxRGBColor

PEXELS_API_KEY = os.getenv("PEXELS_API_KEY", "")


def _search_image(query: str) -> Optional[str]:
    """Search Pexels for a photo. Returns URL or None."""
    if not PEXELS_API_KEY or not query:
        return None
    try:
        resp = requests.get(
            "https://api.pexels.com/v1/search",
            headers={"Authorization": PEXELS_API_KEY},
            params={"query": query, "per_page": 1, "orientation": "landscape"},
            timeout=10,
        )
        data = resp.json()
        photos = data.get("photos", [])
        if photos:
            return photos[0]["src"]["large"]
    except Exception as e:
        print(f"    ⚠️ Image search failed: {e}")
    return None


def _download_image(url: str, dest: Path) -> bool:
    try:
        r = requests.get(url, timeout=15)
        if r.status_code == 200:
            dest.write_bytes(r.content)
            return True
    except Exception as e:
        print(f"    ⚠️ Image download failed: {e}")
    return False


def _hex_to_rgb(hex_str: str) -> RGBColor:
    hex_str = hex_str.lstrip("#")
    if len(hex_str) == 3:
        hex_str = "".join([c * 2 for c in hex_str])
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def _hex_to_pptx_rgb(hex_str: str) -> PptxRGBColor:
    hex_str = hex_str.lstrip("#")
    if len(hex_str) == 3:
        hex_str = "".join([c * 2 for c in hex_str])
    return PptxRGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


# ─── PPTX BUILDER ───────────────────────────────────────────────────────────

def build_pptx(data: Dict, output_path: Path, temp_dir: Path):
    """Build a PowerPoint from the website's JSON schema."""
    print("[builder] Building PowerPoint...")
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    metadata = data.get("metadata", {})
    school_info = metadata.get("school_info", {})
    primary = school_info.get("colors", {}).get("primary", "#8B0000")
    secondary = school_info.get("colors", {}).get("secondary", "#FFD700")
    text_color = "#F5F5F5"
    accent_color = "#63666A"
    bg_color = primary or "#3C0F14"

    slides_data = data.get("slides", [])
    img_dir = temp_dir / "images"
    img_dir.mkdir(exist_ok=True)

    for slide_info in slides_data:
        slide_type = slide_info.get("slide_type", "content")
        content = slide_info.get("content", {})
        layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(layout)

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_pptx_rgb(bg_color)

        # Slide type badge (top left)
        badge = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(0.2),
            PptxInches(3), PptxInches(0.3))
        tf = badge.text_frame
        tf.text = slide_type.replace("_", " ").upper()
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.color.rgb = _hex_to_pptx_rgb(accent_color)

        if slide_type == "title":
            # Title slide
            school_name = content.get("school_name", school_info.get("name", "Your School"))
            unit_title = content.get("unit_title", metadata.get("unit", ""))
            topic_title = content.get("topic_title", metadata.get("topic", ""))
            day_num = content.get("day_number", metadata.get("day_number", 1))
            mascot = school_info.get("mascot", "")

            # School name
            box = slide.shapes.add_textbox(
                PptxInches(0.5), PptxInches(1.5),
                PptxInches(12.333), PptxInches(0.6))
            tf = box.text_frame
            tf.text = school_name
            tf.paragraphs[0].font.size = Pt(28)
            tf.paragraphs[0].font.color.rgb = _hex_to_pptx_rgb(accent_color)
            tf.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER

            # Unit / Day
            box2 = slide.shapes.add_textbox(
                PptxInches(0.5), PptxInches(2.2),
                PptxInches(12.333), PptxInches(0.5))
            tf2 = box2.text_frame
            tf2.text = f"{unit_title} — Day {day_num}"
            tf2.paragraphs[0].font.size = Pt(20)
            tf2.paragraphs[0].font.color.rgb = _hex_to_pptx_rgb(text_color)
            tf2.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER

            # Topic
            box3 = slide.shapes.add_textbox(
                PptxInches(0.5), PptxInches(3.0),
                PptxInches(12.333), PptxInches(1.0))
            tf3 = box3.text_frame
            tf3.text = topic_title
            tf3.paragraphs[0].font.size = Pt(44)
            tf3.paragraphs[0].font.bold = True
            tf3.paragraphs[0].font.color.rgb = _hex_to_pptx_rgb(text_color)
            tf3.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER

            if mascot:
                box4 = slide.shapes.add_textbox(
                    PptxInches(0.5), PptxInches(4.5),
                    PptxInches(12.333), PptxInches(0.4))
                tf4 = box4.text_frame
                tf4.text = mascot
                tf4.paragraphs[0].font.size = Pt(14)
                tf4.paragraphs[0].font.color.rgb = _hex_to_pptx_rgb(accent_color)
                tf4.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER

        else:
            # Content slides
            y = 0.7
            text_width = 9.0

            # Check for image
            has_image = slide_info.get("has_image", False)
            img_query = slide_info.get("image_search_query", "")
            if has_image and img_query:
                text_width = 5.5
                safe_query = re.sub(r'[^\w]', '_', img_query)[:40]
                img_path = img_dir / f"{safe_query}.jpg"
                if not img_path.exists():
                    url = _search_image(img_query)
                    if url:
                        _download_image(url, img_path)
                if img_path.exists():
                    try:
                        slide.shapes.add_picture(
                            str(img_path),
                            PptxInches(7.0), PptxInches(1.0),
                            width=PptxInches(5.5))
                    except Exception as e:
                        print(f"    ⚠️ Failed to add image: {e}")

            # Render content fields in priority order
            field_orders = {
                "hook": ["hook_question", "student_connection"],
                "learning_objective": ["objective_statement", "standard_ref"],
                "journal_prompt": ["prompt_text"],
                "definition_concept": ["term", "definition", "why_it_matters", "real_world_anchor"],
                "real_world_example": ["scenario_title", "scenario_description", "business_name", "numbers_or_data"],
                "comparison": ["concept_a", "concept_b"],
                "activity_intro": ["activity_name", "instructions"],
                "activity_recap": ["connection_to_concept", "recap_points"],
                "practice": ["problem_setup", "problem_1", "problem_2", "hint"],
                "exit_ticket": ["question_1", "question_2", "success_criteria"],
                "next_day_preview": ["preview_text"],
                "prior_review": ["review_question", "answer"],
            }

            # Title for content slide
            title_text = content.get("scenario_title") or content.get("term") or content.get("activity_name") or ""
            if title_text:
                title_box = slide.shapes.add_textbox(
                    PptxInches(0.5), PptxInches(0.4),
                    PptxInches(text_width), PptxInches(0.6))
                tf = title_box.text_frame
                tf.text = title_text
                tf.paragraphs[0].font.size = Pt(32)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = _hex_to_pptx_rgb(text_color)
                y = 1.1

            fields = field_orders.get(slide_type, list(content.keys()))
            for field in fields:
                if field not in content or not content[field]:
                    continue
                val = content[field]

                if field == "recap_points" and isinstance(val, list):
                    val = "\n".join(f"• {p}" for p in val)

                if not val or (isinstance(val, str) and not val.strip()):
                    continue

                # Skip teacher-facing fields that leak to slides
                if field in ("time_estimate", "time_limit", "materials_needed", "grouping"):
                    continue

                box = slide.shapes.add_textbox(
                    PptxInches(0.5), PptxInches(y),
                    PptxInches(text_width), PptxInches(0.8))
                tf = box.text_frame
                tf.word_wrap = True
                tf.text = f"{field.replace('_', ' ').title()}: {val}" if field != "recap_points" else val
                tf.paragraphs[0].font.size = Pt(24)
                tf.paragraphs[0].font.color.rgb = _hex_to_pptx_rgb(text_color)
                y += 0.9

    prs.save(str(output_path))
    print(f"[builder] Saved PowerPoint to {output_path}")


# ─── DOCX BUILDER ─────────────────────────────────────────────────────────────

def build_docx(data: Dict, output_path: Path, doc_type: str = "lesson"):
    """Build a DOCX from the website's JSON schema."""
    print(f"[builder] Building {doc_type} DOCX...")
    doc = Document()

    metadata = data.get("metadata", {})
    school_info = metadata.get("school_info", {})
    plan = data.get("lesson_plan", {})
    activities = data.get("activities", [])

    # Header
    school_name = school_info.get("name", "Your School")
    class_name = metadata.get("class_name", "")
    unit = metadata.get("unit", "")
    day = metadata.get("day_number", 1)
    topic = metadata.get("topic", "Untitled")

    title = doc.add_heading(f"{school_name} — {class_name}", 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    if doc_type == "lesson":
        doc.add_paragraph(f"Unit: {unit} — Day {day}", style="Subtitle").alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_heading(topic, level=1).alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph()

        # Learning Objectives
        objectives = plan.get("learning_objectives", [])
        if objectives:
            doc.add_heading("Learning Objectives", level=2)
            for obj in objectives:
                doc.add_paragraph(obj, style="List Bullet")
            doc.add_paragraph()

        # Duration Breakdown
        breakdown = plan.get("duration_breakdown", [])
        if breakdown:
            doc.add_heading("Duration Breakdown", level=2)
            for item in breakdown:
                phase = item.get("phase", "")
                minutes = item.get("minutes", "")
                desc = item.get("description", "")
                doc.add_paragraph(f"{phase} ({minutes} min): {desc}", style="List Bullet")
            doc.add_paragraph()

        # Instructional Phases
        phases = plan.get("instructional_phases", [])
        if phases:
            doc.add_heading("Instructional Phases", level=2)
            for phase in phases:
                name = phase.get("phase_name", "")
                time = phase.get("time_allocation", "")
                desc = phase.get("description", "")
                doc.add_heading(f"{name} ({time})", level=3)
                doc.add_paragraph(desc)
                notes = phase.get("teacher_notes", "")
                if notes:
                    p = doc.add_paragraph()
                    p.add_run("Teacher Notes: ").bold = True
                    p.add_run(notes)
            doc.add_paragraph()

        # Teacher Notes
        teacher_notes = plan.get("teacher_notes", {})
        if teacher_notes:
            doc.add_heading("Teacher Notes", level=2)
            for key, val in teacher_notes.items():
                if val:
                    p = doc.add_paragraph()
                    p.add_run(f"{key.replace('_', ' ').title()}: ").bold = True
                    p.add_run(str(val))
            doc.add_paragraph()

        # Answer Keys
        answer_keys = plan.get("answer_keys", {})
        if answer_keys:
            doc.add_heading("Answer Keys", level=2)
            for key, val in answer_keys.items():
                if val:
                    p = doc.add_paragraph()
                    p.add_run(f"{key.replace('_', ' ').title()}: ").bold = True
                    p.add_run(str(val))
            doc.add_paragraph()

        # Differentiation
        diff = plan.get("differentiation", {})
        if diff:
            doc.add_heading("Differentiation", level=2)
            for key, val in diff.items():
                if val:
                    p = doc.add_paragraph()
                    p.add_run(f"{key.replace('_', ' ').title()}: ").bold = True
                    p.add_run(str(val))
            doc.add_paragraph()

    elif doc_type == "activity":
        doc.add_heading("Student Activity", level=1).alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph(f"{unit} — Day {day}: {topic}", style="Subtitle").alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph()

        for act in activities:
            doc.add_heading(act.get("activity_name", "Activity"), level=2)
            p = doc.add_paragraph()
            p.add_run("Duration: ").bold = True
            p.add_run(f"{act.get('duration_min', 0)} minutes")
            p = doc.add_paragraph()
            p.add_run("Grouping: ").bold = True
            p.add_run(act.get("grouping", ""))
            p = doc.add_paragraph()
            p.add_run("Materials: ").bold = True
            p.add_run(", ".join(act.get("materials", [])))
            doc.add_paragraph()
            doc.add_heading("Instructions", level=3)
            doc.add_paragraph(act.get("instructions_student_facing", ""))
            doc.add_paragraph()

            deliverables = act.get("deliverables", [])
            if deliverables:
                doc.add_heading("Deliverables", level=3)
                for d in deliverables:
                    doc.add_paragraph(f"{d.get('file_type', '')}: {d.get('description', '')}", style="List Bullet")
                doc.add_paragraph()

    doc.save(str(output_path))
    print(f"[builder] Saved DOCX to {output_path}")
